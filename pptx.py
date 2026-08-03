from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR

# 1. Configuração da Apresentação (Formato Widescreen 16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6]) # Slide em branco

# 2. Cor de Fundo do Slide (Modo Escuro)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(11, 18, 28) # #0B121C

# 3. Funções Auxiliares para criar Caixas e Linhas
def add_box(text, left, top, width, height, bg_color, border_color, is_dashed=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bg_color)
    shape.line.color.rgb = RGBColor(*border_color)
    shape.line.width = Pt(2.5)
    
    if is_dashed:
        shape.line.dash_style = 4 # Estilo tracejado
    
    text_frame = shape.text_frame
    text_frame.clear()
    
    lines = text.split('\n')
    for i, line_str in enumerate(lines):
        p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line_str
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.name = "Arial"
        if i == 0:
            run.font.bold = True
            
    # Limpa parágrafo vazio padrão
    if len(text_frame.paragraphs) > len(lines):
        p = text_frame.paragraphs[0]._element
        p.getparent().remove(p)

    return shape

def add_line(x1, y1, x2, y2, is_dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = RGBColor(165, 165, 165)
    line.line.width = Pt(1.5)
    if is_dashed:
        line.line.dash_style = 4

# 4. Paleta de Cores
bg_main = (11, 18, 28)     # Fundo Trilha Principal
bd_main = (0, 209, 193)    # Borda Cyan (#00D1C1)
bg_exc = (17, 10, 20)      # Fundo Exclusões
bd_exc = (255, 30, 106)    # Borda Magenta (#FF1E6A)
bg_int = (42, 48, 58)      # Fundo Intermediários
bd_int = (165, 165, 165)   # Borda Cinza

# 5. Coordenadas e Dimensões Padrão
w, h = 2.6, 0.9
cx = (13.333 - w) / 2 # Centro horizontal
y0 = 0.2
y_step = 1.2
spacing = 3.8
lx = cx - spacing     # Esquerda
rx = cx + spacing     # Direita

# 6. Criando os Nós (Caixas)
# Nível 1
A = add_box("Participantes Potencialmente Elegíveis\n(n = 969)", cx, y0, w, h, bg_main, bd_main)
Exc1 = add_box("Excluídos (n = 333)\n• Em TTO para TB: n = 6\n• Infecção por MNT: n = 7\n• TB desconhecida: n = 320", cx + 3.2, y0 + 0.4, w+0.2, h+0.4, bg_exc, bd_exc, True)

# Nível 2
B = add_box("Participantes Elegíveis\n(n = 636)", cx, y0 + y_step, w, h, bg_main, bd_main)
Exc2 = add_box("Sem Teste Índice (LF-LAM) (n = 8)\n• CD4 > 200 e sem critério OMS 3/4: n = 7\n• Urina não coletada <24h: n = 1", cx + 3.2, y0 + y_step + 0.4, w+0.2, h+0.4, bg_exc, bd_exc, True)

# Nível 3
C = add_box("Teste Índice (LF-LAM)\n(n = 628)", cx, y0 + 2*y_step, w, h, bg_main, bd_main)

# Nível 4 (Resultados do Teste)
y_d = y0 + 3.1*y_step
D1 = add_box("Teste LF-LAM Negativo\n(n = 443)", lx, y_d, w, h, bg_int, bd_int)
D2 = add_box("Teste LF-LAM Positivo\n(n = 140)", cx, y_d, w, h, bg_int, bd_int)
D3 = add_box("Teste LF-LAM Inconclusivo\n(n = 45)", rx, y_d, w, h, bg_int, bd_int)

# Nível 5 (Referência Padrão)
y_e = y_d + 1.3
E1 = add_box("Com Ref Padrão (MOL/CULT)\n(n = 433)", lx, y_e, w, h, bg_int, bd_int)
E2 = add_box("Com Ref Padrão (MOL/CULT)\n(n = 140)", cx, y_e, w, h, bg_int, bd_int)
E3 = add_box("Com Ref Padrão (MOL/CULT)\n(n = 44)", rx, y_e, w, h, bg_int, bd_int)

# Nível de Exclusões da Referência
Exc3_1 = add_box("Sem Ref Padrão (n = 10)\n• Amostra não coletada: 10", lx - 1.2, y_e - 0.6, w-0.2, h, bg_exc, bd_exc, True)
Exc3_2 = add_box("Sem Ref Padrão\n(n = 0)", cx + 1.5, y_e - 0.6, w-0.4, h-0.2, bg_exc, bd_exc, True)
Exc3_3 = add_box("Sem Ref Padrão (n = 1)\n• Amostra não coletada: 1", rx + 1.6, y_e - 0.6, w-0.4, h, bg_exc, bd_exc, True)

# Nível 6 (Diagnóstico Final)
y_f = y_e + 1.3
F1 = add_box("Diagnóstico Final (M. tb)\n• Presente (+): n = 23\n• Ausente (-): n = 410", lx, y_f, w, h+0.2, bg_main, bd_main)
F2 = add_box("Diagnóstico Final (M. tb)\n• Presente (+): n = 103\n• Ausente (-): n = 37", cx, y_f, w, h+0.2, bg_main, bd_main)
F3 = add_box("Diagnóstico Final (M. tb)\n• Presente (+): n = 7\n• Ausente (-): n = 37", rx, y_f, w, h+0.2, bg_main, bd_main)

# 7. Desenhando as Linhas Conectoras
# Verticais Principais
add_line(cx + w/2, y0 + h, cx + w/2, y0 + y_step)
add_line(cx + w/2, y0 + y_step + h, cx + w/2, y0 + 2*y_step)

# Para Exclusões 1 e 2
add_line(cx + w, y0 + h/2, cx + 3.2, y0 + 0.4 + (h+0.4)/2, True)
add_line(cx + w, y0 + y_step + h/2, cx + 3.2, y0 + y_step + 0.4 + (h+0.4)/2, True)

# Ramificações Teste Índice
add_line(cx + w/2, y0 + 2*y_step + h, lx + w/2, y_d)
add_line(cx + w/2, y0 + 2*y_step + h, cx + w/2, y_d)
add_line(cx + w/2, y0 + 2*y_step + h, rx + w/2, y_d)

# Descendo para Referência Padrão
add_line(lx + w/2, y_d + h, lx + w/2, y_e)
add_line(cx + w/2, y_d + h, cx + w/2, y_e)
add_line(rx + w/2, y_d + h, rx + w/2, y_e)

# Para Exclusões 3
add_line(lx + w/2, y_d + h, lx - 1.2 + (w-0.2)/2, y_e - 0.6, True)
add_line(cx + w, y_d + h/2, cx + 1.5, y_e - 0.6 + (h-0.2)/2, True)
add_line(rx + w, y_d + h/2, rx + 1.6, y_e - 0.6 + h/2, True)

# Descendo para Diagnóstico Final
add_line(lx + w/2, y_e + h, lx + w/2, y_f)
add_line(cx + w/2, y_e + h, cx + w/2, y_f)
add_line(rx + w/2, y_e + h, rx + w/2, y_f)

# 8. Salvar o Arquivo
prs.save('Fluxograma_STARD_Dark.pptx')
print("✅ Arquivo 'Fluxograma_STARD_Dark.pptx' gerado com sucesso!")